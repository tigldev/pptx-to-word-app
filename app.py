import tkinter as tk
import re
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

FONT_CHOICES = [
    "Arial", "Calibri", "Times New Roman", "Courier New",
    "Verdana", "Georgia", "Trebuchet MS", "Comic Sans MS",
    "Lucida Console", "Tahoma"
]

def clean_text(text):
    return re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', text)

def pptx_to_plain_text_word(pptx_path, docx_path, font_name="Arial", font_size=12):
    presentation = Presentation(pptx_path)
    doc = Document()

    for i, slide in enumerate(presentation.slides):
        doc.add_heading(f"Diapositiva {i + 1}", level=1)

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                paragraph = doc.add_paragraph()
                run = paragraph.add_run(clean_text(shape.text))
                font = run.font
                font.name = font_name
                font.size = Pt(font_size)

                rPr = run._element.rPr
                rFonts = OxmlElement('w:rFonts')
                rFonts.set(qn('w:ascii'), font_name)
                rFonts.set(qn('w:hAnsi'), font_name)
                rPr.insert(0, rFonts)

    doc.save(docx_path)

def browse_pptx():
    filepath = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if filepath:
        pptx_entry.delete(0, tk.END)
        pptx_entry.insert(0, filepath)

def save_docx():
    filepath = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word Document", "*.docx")])
    if filepath:
        docx_entry.delete(0, tk.END)
        docx_entry.insert(0, filepath)

def convert():
    pptx_path = pptx_entry.get()
    docx_path = docx_entry.get()
    font_name = font_combobox.get()
    try:
        font_size = int(size_entry.get())
    except ValueError:
        messagebox.showerror("Error", "El tamaño de fuente debe ser un número.")
        return

    if not pptx_path or not docx_path:
        messagebox.showerror("Error", "Debe seleccionar archivos.")
        return

    try:
        pptx_to_plain_text_word(pptx_path, docx_path, font_name, font_size)
        messagebox.showinfo("Éxito", "Conversión completada con éxito.")
    except Exception as e:
        messagebox.showerror("Error", f"Ocurrió un error: {e}")

def actualizar_vista_previa(*args):
    fuente = font_combobox.get()
    try:
        tamaño = int(size_entry.get())
    except ValueError:
        tamaño = 12
    preview_label.config(font=(fuente, tamaño))

# Crear ventana
root = tk.Tk()
root.title("PowerPoint a Word (Texto plano)")
root.configure(bg="#f5f5f5")

# Encabezado
header_label = tk.Label(root, text="PowerPoint a Word App", font=("Segoe UI", 24, "bold"), bg="#f5f5f5", fg="#333")
header_label.grid(row=0, column=0, columnspan=3, pady=(15, 25))

# Campos y botones
tk.Label(root, text="Archivo PowerPoint:", bg="#f5f5f5").grid(row=1, column=0, sticky="e", padx=(15,5), pady=10)
pptx_entry = tk.Entry(root, width=50)
pptx_entry.grid(row=1, column=1, padx=5, pady=10)
tk.Button(root, text="Buscar", command=browse_pptx, bg="#0078d7", fg="white", font=("Segoe UI", 10)).grid(row=1, column=2, padx=(5,15), pady=10)

tk.Label(root, text="Guardar como Word:", bg="#f5f5f5").grid(row=2, column=0, sticky="e", padx=(15,5), pady=10)
docx_entry = tk.Entry(root, width=50)
docx_entry.grid(row=2, column=1, padx=5, pady=10)
tk.Button(root, text="Guardar como", command=save_docx, bg="#0078d7", fg="white", font=("Segoe UI", 10)).grid(row=2, column=2, padx=(5,15), pady=10)

tk.Label(root, text="Fuente:", bg="#f5f5f5").grid(row=3, column=0, sticky="e", padx=(15,5), pady=10)
font_combobox = ttk.Combobox(root, values=FONT_CHOICES, state="readonly", font=("Segoe UI", 10))
font_combobox.set("Arial")
font_combobox.grid(row=3, column=1, sticky="w", padx=5, pady=10)

tk.Label(root, text="Tamaño:", bg="#f5f5f5").grid(row=4, column=0, sticky="e", padx=(15,5), pady=10)
size_entry = tk.Entry(root, font=("Segoe UI", 10))
size_entry.insert(0, "12")
size_entry.grid(row=4, column=1, sticky="w", padx=5, pady=10)

# Conectar eventos para actualizar la vista previa
font_combobox.bind("<<ComboboxSelected>>", actualizar_vista_previa)
size_entry.bind("<KeyRelease>", actualizar_vista_previa)

tk.Button(root, text="Convertir", command=convert, bg="#28a745", fg="white", font=("Segoe UI", 12, "bold")).grid(row=5, column=1, pady=(20,15))

tk.Label(root, text="Vista previa:", bg="#f5f5f5").grid(row=6, column=0, sticky="ne", padx=(15,5), pady=10)
preview_label = tk.Label(root, text="Texto de ejemplo", font=("Arial", 12), bg="#f5f5f5", fg="#333")
preview_label.grid(row=6, column=1, sticky="w", padx=5, pady=10)

root.mainloop()
